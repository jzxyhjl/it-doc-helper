"""
文档内容提取服务
支持PDF、Word、PPT、Markdown、TXT格式
"""
import os
import asyncio
from pathlib import Path
from typing import Optional, Callable
import structlog

logger = structlog.get_logger()


class DocumentExtractor:
    """文档内容提取器"""
    
    # 超时配置
    PDF_EXTRACTION_TIMEOUT = 120  # PDF提取总超时：2分钟
    PDF_PAGE_TIMEOUT = 5  # 单页提取超时：5秒
    MAX_PDF_PAGES = 500  # 最大处理页数（超过则截断）
    MAX_CONTENT_LENGTH = 400000  # 最大内容长度（40万字符，与文档大小验证器保持一致）
    
    @staticmethod
    async def extract_pdf(
        file_path: str, 
        timeout: float = PDF_EXTRACTION_TIMEOUT,
        progress_callback: Optional[Callable[[int, int], None]] = None
    ) -> str:
        """
        提取PDF文档内容（带超时保护）
        
        Args:
            file_path: PDF文件路径
            timeout: 总超时时间（秒）
            progress_callback: 进度回调函数 (current_page, total_pages)
        
        Returns:
            提取的文本内容
        """
        try:
            import pdfplumber
            import concurrent.futures
            
            async def _extract_with_timeout():
                content_parts = []
                total_pages = 0
                processed_pages = 0
                
                # 在线程池中执行同步的PDF操作
                loop = asyncio.get_event_loop()
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    # 先获取总页数
                    def get_page_count():
                        with pdfplumber.open(file_path) as pdf:
                            return len(pdf.pages)
                    
                    total_pages = await loop.run_in_executor(executor, get_page_count)
                    
                    if total_pages > DocumentExtractor.MAX_PDF_PAGES:
                        logger.warning(
                            "PDF页数过多，将截断处理",
                            file_path=file_path,
                            total_pages=total_pages,
                            max_pages=DocumentExtractor.MAX_PDF_PAGES
                        )
                        total_pages = DocumentExtractor.MAX_PDF_PAGES
                    
                    # 分页提取，每页设置超时
                    with pdfplumber.open(file_path) as pdf:
                        for page_num, page in enumerate(pdf.pages[:total_pages], 1):
                            # 检查总内容长度
                            current_length = sum(len(part) for part in content_parts)
                            if current_length > DocumentExtractor.MAX_CONTENT_LENGTH:
                                logger.warning(
                                    "PDF内容过长，提前截断",
                                    file_path=file_path,
                                    current_length=current_length,
                                    max_length=DocumentExtractor.MAX_CONTENT_LENGTH,
                                    processed_pages=processed_pages
                                )
                                break
                            
                            # 单页提取（带超时）
                            def extract_page_text(p):
                                return p.extract_text() or ""
                            
                            try:
                                page_text = await asyncio.wait_for(
                                    loop.run_in_executor(executor, extract_page_text, page),
                                    timeout=DocumentExtractor.PDF_PAGE_TIMEOUT
                                )
                                
                                if page_text:
                                    content_parts.append(page_text)
                                
                                processed_pages += 1
                                
                                # 进度回调（异步）
                                if progress_callback:
                                    try:
                                        await progress_callback(processed_pages, total_pages)
                                    except Exception:
                                        pass  # 忽略回调错误
                                
                            except asyncio.TimeoutError:
                                logger.warning(
                                    f"PDF第{page_num}页提取超时，跳过",
                                    file_path=file_path,
                                    page_num=page_num,
                                    timeout=DocumentExtractor.PDF_PAGE_TIMEOUT
                                )
                                # 跳过超时的页面，继续处理下一页
                                processed_pages += 1
                                continue
                            except Exception as e:
                                logger.warning(
                                    f"PDF第{page_num}页提取失败，跳过",
                                    file_path=file_path,
                                    page_num=page_num,
                                    error=str(e)
                                )
                                processed_pages += 1
                                continue
                
                content = "\n\n".join(content_parts)
                logger.info(
                    "PDF内容提取成功",
                    file_path=file_path,
                    total_pages=total_pages,
                    processed_pages=processed_pages,
                    content_length=len(content)
                )
                return content
            
            # 执行提取，带总超时保护
            content = await asyncio.wait_for(_extract_with_timeout(), timeout=timeout)
            return content
            
        except asyncio.TimeoutError:
            logger.error(
                "PDF提取总超时",
                file_path=file_path,
                timeout=timeout
            )
            raise Exception(f"PDF内容提取超时（超过{timeout}秒），文件可能过大或损坏")
        except Exception as e:
            logger.error("PDF提取失败", file_path=file_path, error=str(e))
            raise Exception(f"PDF内容提取失败: {str(e)}")
    
    @staticmethod
    async def extract_word(file_path: str) -> str:
        """提取Word文档内容"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            content_parts = []
            
            # 提取段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)
            
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        content_parts.append(row_text)
            
            content = "\n".join(content_parts)
            logger.info("Word内容提取成功", file_path=file_path)
            return content
            
        except Exception as e:
            logger.error("Word提取失败", file_path=file_path, error=str(e))
            raise Exception(f"Word内容提取失败: {str(e)}")
    
    @staticmethod
    async def extract_ppt(file_path: str) -> str:
        """提取PPT文档内容"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                
                # 提取幻灯片标题和内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    content_parts.append(f"--- 幻灯片 {slide_num} ---")
                    content_parts.extend(slide_content)
                    content_parts.append("")  # 空行分隔
            
            content = "\n".join(content_parts)
            logger.info("PPT内容提取成功", file_path=file_path, slides=len(prs.slides))
            return content
            
        except Exception as e:
            logger.error("PPT提取失败", file_path=file_path, error=str(e))
            raise Exception(f"PPT内容提取失败: {str(e)}")
    
    @staticmethod
    async def extract_markdown(file_path: str) -> str:
        """提取Markdown文档内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            logger.info("Markdown内容提取成功", file_path=file_path)
            return content
            
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                import chardet
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    encoding = chardet.detect(raw_data)['encoding']
                
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                logger.info("Markdown内容提取成功（检测编码）", file_path=file_path, encoding=encoding)
                return content
            except Exception as e:
                logger.error("Markdown提取失败", file_path=file_path, error=str(e))
                raise Exception(f"Markdown内容提取失败: {str(e)}")
        except Exception as e:
            logger.error("Markdown提取失败", file_path=file_path, error=str(e))
            raise Exception(f"Markdown内容提取失败: {str(e)}")
    
    @staticmethod
    async def extract_txt(file_path: str) -> str:
        """提取TXT文档内容"""
        try:
            import chardet
            
            # 检测编码
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # 读取内容
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            logger.info("TXT内容提取成功", file_path=file_path, encoding=encoding)
            return content
            
        except Exception as e:
            logger.error("TXT提取失败", file_path=file_path, error=str(e))
            raise Exception(f"TXT内容提取失败: {str(e)}")
    
    @staticmethod
    async def extract(
        file_path: str, 
        file_type: str,
        timeout: Optional[float] = None,
        progress_callback: Optional[Callable[[int, int], None]] = None
    ) -> str:
        """
        根据文件类型提取内容（带超时保护）
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（pdf/docx/pptx/md/txt）
            timeout: 提取超时时间（秒），None则使用默认值
            progress_callback: 进度回调函数（仅PDF支持）
        
        Returns:
            提取的文本内容
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_type = file_type.lower()
        
        # 为不同文件类型设置默认超时
        default_timeouts = {
            'pdf': DocumentExtractor.PDF_EXTRACTION_TIMEOUT,
            'docx': 60,  # Word: 1分钟
            'pptx': 60,  # PPT: 1分钟
            'md': 10,    # Markdown: 10秒
            'markdown': 10,
            'txt': 10,   # TXT: 10秒
        }
        
        extraction_timeout = timeout or default_timeouts.get(file_type, 60)
        
        # 为所有提取操作添加超时保护
        async def extract_with_timeout():
            if file_type == 'pdf':
                return await DocumentExtractor.extract_pdf(
                    file_path, 
                    timeout=extraction_timeout,
                    progress_callback=progress_callback
                )
            elif file_type == 'docx':
                return await DocumentExtractor.extract_word(file_path)
            elif file_type == 'pptx':
                return await DocumentExtractor.extract_ppt(file_path)
            elif file_type in ['md', 'markdown']:
                return await DocumentExtractor.extract_markdown(file_path)
            elif file_type == 'txt':
                return await DocumentExtractor.extract_txt(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_type}")
        
        try:
            content = await asyncio.wait_for(
                extract_with_timeout(),
                timeout=extraction_timeout
            )
            return content
        except asyncio.TimeoutError:
            logger.error(
                "文档提取超时",
                file_path=file_path,
                file_type=file_type,
                timeout=extraction_timeout
            )
            raise Exception(
                f"{file_type.upper()}内容提取超时（超过{extraction_timeout}秒），"
                "文件可能过大或损坏。建议拆分后处理。"
            )

